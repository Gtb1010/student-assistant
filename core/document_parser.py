"""
core/document_parser.py — Nxjerr tekst nga PDF, DOCX, PPTX.

LlamaIndex përdoret VETËM për strukturën Document (text + metadata).
Gjithë logjika e parsimit është e shkruar manualisht.
"""

from pathlib import Path
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from llama_index.core import Document   # e vetmja gjë e LlamaIndex-it


class DocumentParser:
    """Kthen listë Document-esh nga skedari i dhënë."""

    def parse(self, path: str | Path) -> list[Document]:
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"Skedari nuk u gjet: {path}")
        ext = path.suffix.lower()
        if ext == ".pdf":   return self._pdf(path)
        if ext == ".docx":  return self._docx(path)
        if ext == ".pptx":  return self._pptx(path)
        raise ValueError(f"Lloji '{ext}' nuk mbështetet. Lejuar: .pdf .docx .pptx")

    @staticmethod
    def _pdf(path: Path) -> list[Document]:
        reader = PdfReader(str(path))
        docs = []
        for i, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                docs.append(Document(text=text, metadata={
                    "source": path.name, "file_path": str(path),
                    "file_type": "pdf", "page_number": i,
                    "page_count": len(reader.pages),
                }))
        return docs

    @staticmethod
    def _docx(path: Path) -> list[Document]:
        lines = [p.text for p in DocxDocument(str(path)).paragraphs if p.text.strip()]
        text  = "\n".join(lines)
        return [Document(text=text, metadata={
            "source": path.name, "file_path": str(path), "file_type": "docx",
        })] if text.strip() else []

    @staticmethod
    def _pptx(path: Path) -> list[Document]:
        prs, docs = Presentation(str(path)), []
        for i, slide in enumerate(prs.slides, 1):
            lines = [
                p.text.strip()
                for shape in slide.shapes if shape.has_text_frame
                for p in shape.text_frame.paragraphs if p.text.strip()
            ]
            if lines:
                docs.append(Document(text="\n".join(lines), metadata={
                    "source": path.name, "file_path": str(path),
                    "file_type": "pptx", "slide_number": i,
                    "slide_count": len(prs.slides),
                }))
        return docs